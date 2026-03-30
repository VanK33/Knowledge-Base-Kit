#!/usr/bin/env python3
"""
Extract text content from various file types for content matching.

Usage:
  python3 extract_text.py <file_path>
  echo '<file_path>' | python3 extract_text.py

Output: extracted text to stdout (UTF-8). Exit code 0 on success.

Supported formats:
  .docx  — requires python-docx (pip install python-docx)
  .xlsx  — requires openpyxl (pip install openpyxl)
  .pptx  — requires python-pptx (pip install python-pptx)
  .csv   — stdlib csv module
  .txt   — direct read
  .md    — direct read
  .json  — direct read
  .pdf   — first 5 pages text (requires PyMuPDF: pip install PyMuPDF)

For .png/.jpg/.jpeg — use Claude's Read tool (OCR) instead of this script.
"""

import csv
import io
import json
import sys
from pathlib import Path


def extract_docx(path: Path) -> str:
    """Extract text from .docx file."""
    try:
        from docx import Document
    except ImportError:
        return f"[extract_text] python-docx not installed. Run: pip install python-docx"

    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    """Extract text from .xlsx file."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return f"[extract_text] openpyxl not installed. Run: pip install openpyxl"

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    """Extract text from .pptx file."""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[extract_text] python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide_texts:
            parts.append(f"--- Slide {i} ---")
            parts.extend(slide_texts)
    return "\n".join(parts)


def extract_csv_file(path: Path) -> str:
    """Extract text from .csv file."""
    parts = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i >= 200:  # limit rows
                parts.append(f"... ({i}+ rows)")
                break
            cells = [c.strip() for c in row if c.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    """Extract text from first pages of PDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return f"[extract_text] PyMuPDF not installed. Run: pip install PyMuPDF"

    doc = fitz.open(str(path))
    parts = []
    max_pages = min(5, len(doc))
    for i in range(max_pages):
        text = doc[i].get_text().strip()
        if text:
            parts.append(f"--- Page {i + 1} ---")
            parts.append(text)
    doc.close()
    return "\n".join(parts)


def extract_text_file(path: Path) -> str:
    """Read plain text files directly."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read(50000)  # limit to ~50KB


def extract_json_file(path: Path) -> str:
    """Extract text representation from JSON."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)[:50000]


EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
    ".pptx": extract_pptx,
    ".csv": extract_csv_file,
    ".tsv": extract_csv_file,
    ".pdf": extract_pdf,
    ".txt": extract_text_file,
    ".md": extract_text_file,
    ".json": extract_json_file,
    ".yaml": extract_text_file,
    ".yml": extract_text_file,
    ".xml": extract_text_file,
    ".html": extract_text_file,
    ".htm": extract_text_file,
}


def extract(file_path: str) -> str:
    """Extract text from a file. Returns extracted text or error message."""
    path = Path(file_path)

    if not path.exists():
        return f"[extract_text] File not found: {path}"

    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext)

    if extractor is None:
        return f"[extract_text] Unsupported format: {ext}. Supported: {', '.join(sorted(EXTRACTORS.keys()))}"

    try:
        return extractor(path)
    except Exception as e:
        return f"[extract_text] Error extracting {ext}: {e}"


def main():
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
    else:
        file_path = sys.stdin.read().strip()

    if not file_path:
        print("[extract_text] No file path provided", file=sys.stderr)
        sys.exit(1)

    result = extract(file_path)
    print(result)


if __name__ == "__main__":
    main()
